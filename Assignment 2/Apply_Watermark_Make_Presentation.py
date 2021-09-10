# Link: https://stackoverflow.com/questions/44275443/python-inserts-pictures-to-powerpoint-how-to-set-the-width-and-height-of-the-pi
# Import library
import os
from wand.image import Image
import pptx
import pytest

#make a list with all the jpg images
def create_image_list():
    list_images = []
    for f in os.listdir('Images'):
        if os.path.splitext(f)[1].lower() in ('.jpg', '.jpeg'):
            list_images.append('Images/'+f)
    return list_images

#apply logo on images
def make_watermark_images(image_list):
    watermark_images_list = []
    for i in image_list:
        # Import the image
        with Image(filename = i) as image:
            # Import the watermark image
            with Image(filename ='Images/nike_black.png') as logo:
                logo.resize(width=1600, height=600)
                # Clone the image in order to process
                with image.clone() as watermark_image:
                    # Invoke watermark function with watermark image, transparency as 0.5, left as 10 and top as 20
                    watermark_image.watermark(logo, 0.1, 20,20)
                    # Save the image
                    output_path = 'Images/Watermark images'
                    if not os.path.exists(output_path):
                        os.mkdir(output_path)
                    watermark_image.save(filename ='Images/Watermark images/watermark'+str(image_list.index(i))+'.jpg')
                    watermark_images_list.append('Images/Watermark images/watermark'+str(image_list.index(i))+'.jpg')
    return watermark_images_list

def _add_image(slide, image_url, slide_width, slide_height):

    # Calculate the image size of the image
    im = Image(filename = image_url)
    width, height = im.size

    # Calculate ratios and compare
    image_ratio = width / height
    slide_ratio = slide_width/slide_height
    ratio_difference = image_ratio/slide_ratio
    # print(ratio_difference)
    # print(slide_width*ratio_difference,slide_height*ratio_difference)

    # Insert the picture; applied some mathematics
    # pic   = slide.shapes.add_picture(image_url,int(slide_width*0.1),int(slide_height*0.3),Inches(1))
    pic   = slide.shapes.add_picture(image_url,int(slide_width*0.1),int(slide_height*0.3),int(slide_width*0.3),int(slide_width*(height/width)*0.3))
   
    return "successful"

# function for creating ppt
def create_presentation(images_list):
    OUTPUT_TAG = "MY PRESENTATION"

    # new
    prs = pptx.Presentation()

    # title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    # set title
    title = slide.shapes.title
    title.text = OUTPUT_TAG
    i = 1

    for g in images_list:
        # print(g)
        # fix layout: Heading, Subheading, image
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # title for slide
        title = slide.shapes.title
        title.text = "Sample Presentation"
        title.text_frame.paragraphs[0].font.size = pptx.util.Pt(24)

        # subheading text
        sub = slide.placeholders[1].text = "Sample Subtitle: "+str(i)
        slide.placeholders[1].text_frame.paragraphs[0].font.size = pptx.util.Pt(18)

        # function for adding image in slide
        _add_image(slide,g,prs.slide_width,prs.slide_height)
        i=i+1

    # save and load presentation
    prs.save("%s.pptx" % OUTPUT_TAG)
    os.startfile("%s.pptx" % OUTPUT_TAG)

# create image list on which watermark need to be applied
list_images = create_image_list()
# print(list_images)

# Get all watermarked images
list_watermark_images = make_watermark_images(list_images)
# print(list_watermark_images)

#create presentation
create_presentation(list_watermark_images)

# Test Case 1:
def test_one():
    # Fail
    assert make_watermark_images(['Images/image1.jpg', 'Images/image2.jpg'])  == ['Images/Watermark images/watermark5.jpg', 'Images/Watermark images/watermark1.jpg']
    # Pass
    # assert make_watermark_images(['Images/image1.jpg', 'Images/image2.jpg'])  == ['Images/Watermark images/watermark0.jpg', 'Images/Watermark images/watermark1.jpg']

# Test Case 2:
def test_two():
    # Pass
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    g = 'Images/Watermark images/watermark0.jpg'
    assert _add_image(slide,g,9144000,5143500) == "successful"
    # Fail
    # assert _add_image(slide,g,9144000,5143500) == "Failed"
    